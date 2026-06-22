"""Post-generation validation — check all shapes are within slide bounds"""
from pptx import Presentation

def validate(path):
    prs = Presentation(path)
    issues = []
    pw, ph = prs.slide_width, prs.slide_height
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            try:
                left = shape.left
                top = shape.top
                right = left + shape.width
                bottom = top + shape.height
            except Exception:
                continue
            name = getattr(shape, 'name', f'shape_{i}')
            if left < -1:
                issues.append(f"Slide {i+1}: '{name}' left={left} < 0")
            if top < -1:
                issues.append(f"Slide {i+1}: '{name}' top={top} < 0")
            if right > pw + 1:
                issues.append(f"Slide {i+1}: '{name}' overflow right ({right} > {pw})")
            if bottom > ph + 1:
                issues.append(f"Slide {i+1}: '{name}' overflow bottom ({bottom} > {ph})")
    if issues:
        print(f"❌ {len(issues)} validation issues:")
        for issue in issues:
            print(f"  - {issue}")
        return False
    print(f"✅ All {len(prs.slides)} slides validated ({pw}x{ph})")
    return True
