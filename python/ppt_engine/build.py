"""PPT build pipeline — theme → layout → render → animation → validate"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from pptx import Presentation
from pptx.util import Inches
from ppt_engine.theme import Theme
from ppt_engine.layout import LayoutGrid
from ppt_engine.renderer import render_cover, render_divider, render_card_grid
from ppt_engine.animation import apply_transition
from ppt_engine.validate import validate

def build_pptx(output_path, title, pages, theme_name="default"):
    """Main entry: takes structured page data, returns pptx path.

    pages = [
        {'type': 'cover',     'title': '...', 'subtitle': '...'},
        {'type': 'divider',   'title': 'Section 1'},
        {'type': 'card_grid', 'rows': 2, 'cols': 2, 'cards': [
            {'title': 'Card 1', 'bullets': ['item a', 'item b'], 'img': None},
            {'title': 'Card 2', 'bullets': ['item c'], 'img_url': 'https://...'},
            ...
        ]},
    ]
    Cards accept: img (local path), img_url (HTTP URL) — both auto-preprocessed.
    """
    theme = Theme(theme_name)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for page in pages:
        ptype = page.get('type', 'card_grid')
        if ptype == 'cover':
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            render_cover(slide, theme, page['title'], page.get('subtitle', ''))
            apply_transition(slide, 'cover')

        elif ptype == 'divider':
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            render_divider(slide, theme, page['title'])
            apply_transition(slide, 'divider')

        elif ptype == 'card_grid':
            grid = LayoutGrid()
            rows, cols = page.get('rows', 2), page.get('cols', 2)
            cells = grid.grid(rows, cols, page.get('row_ratios'), page.get('col_ratios'))
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            render_card_grid(slide, theme, cells, page.get('cards', []))
            apply_transition(slide, page.get('transition', 'content'))

    prs.save(output_path)
    validate(output_path)
    return output_path


# ── Demo ──
if __name__ == '__main__':
    demo_pages = [
        {'type': 'cover', 'title': 'AI-Powered PPT Engine', 'subtitle': 'Layout Math + Image Pipeline + Auto-Validation'},
        {'type': 'divider', 'title': 'Core Architecture'},
        {'type': 'card_grid', 'rows': 2, 'cols': 2, 'cards': [
            {'title': 'Image Pipeline', 'bullets': ['Center-crop to target ratio', 'Zero stretching or distortion', 'JPEG quality 85 compression']},
            {'title': 'Layout Engine', 'bullets': ['Constraint-based grid system', 'Adaptive card heights', 'Built-in overflow asserts']},
            {'title': 'Renderer', 'bullets': ['Separation of data & rendering', 'Theme-swappable color system', 'Single-line card creation']},
            {'title': 'Validation', 'bullets': ['Post-generation bounds check', 'All shapes verified in-page', 'Immediate error reporting']},
        ]},
    ]
    build_pptx('/tmp/demo_output.pptx', 'Demo PPT', demo_pages)
    print("Done → /tmp/demo_output.pptx")
