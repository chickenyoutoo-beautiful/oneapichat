"""Renderer — draws computed layouts onto python-pptx slides"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from .image_processor import ImageProcessor
from .layout import CardLayout

_ip = ImageProcessor()

def _add_bg(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def _add_textbox(slide, x, y, w, h, text, size=10, color=None, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.alignment = align
    if color:
        p.font.color.rgb = color
    return tf

class CardRenderer:
    def __init__(self, slide, theme):
        self.slide = slide
        self.theme = theme

    def render(self, layout):
        r = layout
        # Card background
        _add_bg(self.slide, r['card_x'], r['card_y'], r['card_w'], r['card_h'], self.theme.card_bg)

        # Image
        if 'img' in r and r.get('img_path'):
            x, y, w, h = r['img']
            processed = _ip.process(r['img_path'], w, h)
            if processed:
                self.slide.shapes.add_picture(processed, Inches(x), Inches(y), Inches(w), Inches(h))

        # Title
        x, y, w, h = r['title']
        _add_textbox(self.slide, x, y, w, h, r.get('title_text', ''),
                     size=12.5, color=self.theme.accent, bold=True)

        # Separator
        x, y, w, h = r['sep']
        _add_bg(self.slide, x, y, w, h, self.theme.accent)

        # Bullets
        x, y, w, h = r['bullets']
        tf = _add_textbox(self.slide, x, y, w, h + 0.15, '')
        for i, line in enumerate(r.get('bullet_lines', [])):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = '▸ ' + line
            p.font.size = Pt(9.5)
            p.font.color.rgb = self.theme.text
            p.space_after = Pt(2)


def render_cover(slide, theme, title, subtitle="", bg_color=None):
    bg = bg_color or theme.accent
    _add_bg(slide, 0, 0, 10, 7.5, bg)
    _add_textbox(slide, 0.8, 2.5, 8.4, 1.2, title, size=32,
                 color=theme.text_inverse, bold=True, align=PP_ALIGN.LEFT)
    if subtitle:
        _add_textbox(slide, 0.8, 3.8, 8.4, 0.6, subtitle, size=14,
                     color=RGBColor(0xCC, 0xCC, 0xEE))


def render_divider(slide, theme, title):
    _add_bg(slide, 0, 0, 10, 7.5, theme.dark_bg)
    _add_textbox(slide, 0.8, 3.2, 8.4, 0.8, title, size=24,
                 color=theme.text_inverse, bold=True, align=PP_ALIGN.CENTER)


def render_card_grid(slide, theme, grid_cells, cards_data):
    """cards_data: list of {title, bullets, img} matching grid_cells order"""
    renderer = CardRenderer(slide, theme)
    layout = CardLayout()
    for i, cell in enumerate(grid_cells):
        if i >= len(cards_data):
            break
        data = cards_data[i]
        x, y, w, h = cell
        result = layout.compute(x, y, w, h, data.get('title', ''),
                                data.get('bullets', []), data.get('img'))
        renderer.render(result)
